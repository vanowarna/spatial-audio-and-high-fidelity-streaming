"""
Build the final 5-minute Spatial Audio presentation.

Style: Neumorphism — soft off-white surface, dual-direction shadows on cards,
       rounded corners, Google Sans / Roboto typography, Google color accents.
"""

import os
from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree

REPO_ROOT = Path("/sessions/busy-zen-hamilton/mnt/spatial-audio-and-high-fidelity-streaming")
FIG_DIR = REPO_ROOT / "report" / "figures"
OUT_DIR = REPO_ROOT / "deliverables"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PPTX = OUT_DIR / "Spatial_Audio_Presentation.pptx"

# --- Colors (neumorphic surface + Google accents) ---
SURFACE = RGBColor(0xE6, 0xEA, 0xF0)      # neumorphic base
SURFACE_HEX = "E6EAF0"
SHADOW_LIGHT_HEX = "FFFFFF"               # top-left highlight
SHADOW_DARK_HEX = "BCC4D0"                # bottom-right shadow
TEXT_PRIMARY = RGBColor(0x1F, 0x24, 0x33)
TEXT_SECONDARY = RGBColor(0x52, 0x5B, 0x6E)
TEXT_DIM = RGBColor(0x8C, 0x96, 0xA8)

GOOGLE_BLUE = RGBColor(0x42, 0x85, 0xF4)
GOOGLE_BLUE_HEX = "4285F4"
GOOGLE_RED = RGBColor(0xEA, 0x43, 0x35)
GOOGLE_RED_HEX = "EA4335"
GOOGLE_YELLOW = RGBColor(0xFB, 0xBC, 0x04)
GOOGLE_YELLOW_HEX = "FBBC04"
GOOGLE_GREEN = RGBColor(0x34, 0xA8, 0x53)
GOOGLE_GREEN_HEX = "34A853"

FONT_HEAD = "Google Sans"     # falls back to Roboto / system sans on most installs
FONT_BODY = "Roboto"


# ---------- low-level helpers ----------

def _ns():
    return {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

def _qa(tag):
    return '{http://schemas.openxmlformats.org/drawingml/2006/main}' + tag


def _add_outer_shadow(shape, dist=12700, blur=63500, alpha=35000, direction=2700000, color="BCC4D0"):
    """Apply a soft outer drop shadow to a shape (neumorphic dark side)."""
    sppr = shape._element.find(qn("p:spPr"))
    # remove existing effectLst
    for child in sppr.findall(qn("a:effectLst")):
        sppr.remove(child)
    eff = etree.SubElement(sppr, qn("a:effectLst"))
    outer = etree.SubElement(eff, qn("a:outerShdw"))
    outer.set("blurRad", str(blur))
    outer.set("dist", str(dist))
    outer.set("dir", str(direction))
    outer.set("rotWithShape", "0")
    srgb = etree.SubElement(outer, qn("a:srgbClr"))
    srgb.set("val", color)
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(alpha))


def _add_inner_shadow(shape, dist=8000, blur=40000, alpha=45000, direction=2700000, color="BCC4D0"):
    """Optional inset shadow for pressed-in look."""
    sppr = shape._element.find(qn("p:spPr"))
    for child in sppr.findall(qn("a:effectLst")):
        sppr.remove(child)
    eff = etree.SubElement(sppr, qn("a:effectLst"))
    inner = etree.SubElement(eff, qn("a:innerShdw"))
    inner.set("blurRad", str(blur))
    inner.set("dist", str(dist))
    inner.set("dir", str(direction))
    srgb = etree.SubElement(inner, qn("a:srgbClr"))
    srgb.set("val", color)
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(alpha))


def _set_no_outline(shape):
    sppr = shape._element.find(qn("p:spPr"))
    for ln in sppr.findall(qn("a:ln")):
        sppr.remove(ln)
    ln = etree.SubElement(sppr, qn("a:ln"))
    nf = etree.SubElement(ln, qn("a:noFill"))


def neumorphic_card(slide, x, y, w, h, *, fill_hex=SURFACE_HEX, raised=True, accent_hex=None):
    """Add a neumorphic card with soft dual shadows."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.10
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    _set_no_outline(shape)
    if raised:
        # dark shadow lower-right
        _add_outer_shadow(shape, dist=18000, blur=80000, alpha=42000,
                          direction=2700000, color=SHADOW_DARK_HEX)
    else:
        # pressed
        _add_inner_shadow(shape, dist=10000, blur=40000, alpha=45000,
                          direction=2700000, color=SHADOW_DARK_HEX)
    if accent_hex:
        # decorative left bar
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Emu(0), Emu(70000), h)
        bar.adjustments[0] = 0.45
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
        _set_no_outline(bar)
        _add_outer_shadow(bar, dist=4000, blur=20000, alpha=20000, color=accent_hex)
    return shape


def add_text(slide, text, x, y, w, h, *, font=FONT_BODY, size=14, bold=False,
             color=TEXT_PRIMARY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, space_after_pt=2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    first = True
    for line in text:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after_pt)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_pill(slide, text, x, y, *, fill_hex=GOOGLE_BLUE_HEX, text_color=RGBColor(0xFF, 0xFF, 0xFF),
             padding_x=Emu(120000), padding_y=Emu(30000), size=11, bold=True):
    # Approximate width based on text length and font size
    # ~ pt size * 5500 EMU per character is a good estimate for sans bold
    char_w = int(size * 7500)
    w = max(Emu(700000), padding_x * 2 + char_w * len(text))
    h = Emu(380000)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    _set_no_outline(pill)
    _add_outer_shadow(pill, dist=8000, blur=30000, alpha=30000, color=fill_hex)
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_HEAD
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return pill


def add_circle(slide, x, y, d, color_hex):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    _set_no_outline(c)
    _add_outer_shadow(c, dist=6000, blur=25000, alpha=25000, color=color_hex)
    return c


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def style_slide_background(slide, color_hex=SURFACE_HEX):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_hex)


def add_footer(slide, slide_num, total, *, label="Spatial Audio & High-Fidelity Streaming"):
    add_text(slide, label, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             font=FONT_BODY, size=10, color=TEXT_DIM, align=PP_ALIGN.LEFT)
    add_text(slide, f"{slide_num} / {total}", Inches(11.7), Inches(7.05),
             Inches(1.7), Inches(0.3),
             font=FONT_BODY, size=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


# ---------- slide builders ----------

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    TOTAL = 10

    # ===== SLIDE 1 — TITLE =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s)
    # Big neumorphic central card
    neumorphic_card(s, Inches(1.3), Inches(1.4), Inches(10.7), Inches(4.7),
                    fill_hex="EEF1F6")
    # Decorative Google-color dots
    add_circle(s, Inches(1.7), Inches(1.8), Inches(0.4), GOOGLE_BLUE_HEX)
    add_circle(s, Inches(2.2), Inches(1.8), Inches(0.4), GOOGLE_RED_HEX)
    add_circle(s, Inches(2.7), Inches(1.8), Inches(0.4), GOOGLE_YELLOW_HEX)
    add_circle(s, Inches(3.2), Inches(1.8), Inches(0.4), GOOGLE_GREEN_HEX)
    # Tag
    add_text(s, "PROJECT 4 · 2102571 MULTIMEDIA COMMUNICATION",
             Inches(1.7), Inches(2.4), Inches(10), Inches(0.4),
             font=FONT_HEAD, size=12, bold=True, color=GOOGLE_BLUE)
    # Title
    add_text(s, "Spatial Audio &", Inches(1.7), Inches(2.85), Inches(10), Inches(0.9),
             font=FONT_HEAD, size=44, bold=True, color=TEXT_PRIMARY)
    add_text(s, "High-Fidelity Streaming",
             Inches(1.7), Inches(3.7), Inches(11), Inches(0.9),
             font=FONT_HEAD, size=44, bold=True, color=TEXT_PRIMARY)
    add_text(s, "Evaluating Opus, AAC and MP3 for low-latency binaural audio over constrained networks",
             Inches(1.7), Inches(4.7), Inches(10), Inches(0.6),
             font=FONT_BODY, size=14, italic=True, color=TEXT_SECONDARY)

    # Authors row
    add_text(s, "Yuwadee Tongkong  ·  Vanodhya Warnasooriya  ·  Chanakan Hambleton",
             Inches(1.7), Inches(5.4), Inches(10), Inches(0.4),
             font=FONT_BODY, size=12, color=TEXT_SECONDARY)
    add_text(s, "Department of Electrical Engineering · Chulalongkorn University · May 2026",
             Inches(1.7), Inches(5.75), Inches(10), Inches(0.4),
             font=FONT_BODY, size=10, color=TEXT_DIM)
    add_speaker_notes(s,
        "Good morning. We are presenting Project 4: Spatial Audio and High-Fidelity Streaming. "
        "We compare three perceptual codecs - MP3, AAC and Opus - for low-latency binaural delivery "
        "and demonstrate a working spatial audio web player. I'll keep us to five minutes. "
        "[10 seconds]"
    )
    add_footer(s, 1, TOTAL)

    # ===== SLIDE 2 — Why Spatial Audio? =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s)
    # Title
    add_text(s, "Why Spatial Audio?", Inches(0.6), Inches(0.4), Inches(11), Inches(0.7),
             font=FONT_HEAD, size=32, bold=True, color=TEXT_PRIMARY)
    add_text(s, "The bandwidth gap that quietly breaks immersion",
             Inches(0.6), Inches(1.05), Inches(11), Inches(0.4),
             font=FONT_BODY, size=13, italic=True, color=TEXT_SECONDARY)

    # Three cards
    cards = [
        ("Bandwidth gap", "Video uses 3–8 Mbps;\naudio is capped at 64–128 kbps.",
         GOOGLE_BLUE_HEX),
        ("Spatial cues are fragile", "ITD and ILD live in transients and\nhigh frequencies — the first to be cut.",
         GOOGLE_RED_HEX),
        ("Immersion at risk", "VR / AR / Metaverse demand\npresence — bad audio breaks it.",
         GOOGLE_YELLOW_HEX),
    ]
    for i, (h, body, color) in enumerate(cards):
        x = Inches(0.6 + i * 4.15)
        neumorphic_card(s, x, Inches(1.7), Inches(3.95), Inches(2.6))
        add_circle(s, x + Inches(0.3), Inches(1.95), Inches(0.55), color)
        add_text(s, h, x + Inches(0.3), Inches(2.65), Inches(3.5), Inches(0.5),
                 font=FONT_HEAD, size=18, bold=True, color=TEXT_PRIMARY)
        add_text(s, body, x + Inches(0.3), Inches(3.15), Inches(3.5), Inches(1.2),
                 font=FONT_BODY, size=13, color=TEXT_SECONDARY)

    # Bottom callout
    neumorphic_card(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.0),
                    fill_hex="EFF4FB", accent_hex=GOOGLE_GREEN_HEX)
    add_text(s, "Question",
             Inches(1.0), Inches(4.85), Inches(11), Inches(0.5),
             font=FONT_HEAD, size=12, bold=True, color=GOOGLE_GREEN)
    add_text(s,
             "Which codec converts the smallest bitrate into the most directional clarity?",
             Inches(1.0), Inches(5.25), Inches(11), Inches(1.4),
             font=FONT_HEAD, size=22, bold=True, color=TEXT_PRIMARY)
    add_speaker_notes(s,
        "Modern streaming spends most of its budget on video. Audio gets the leftover. "
        "That's a problem for spatial audio because the cues that tell our brain where a sound is "
        "coming from - the inter-aural time and level differences - sit precisely in the high-frequency "
        "and transient detail that aggressive compression discards first. "
        "So our central question is simple: which codec converts the smallest bitrate into the most "
        "directional clarity? [35 seconds]"
    )
    add_footer(s, 2, TOTAL)

    # ===== SLIDE 3 — Objectives =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s)
    add_text(s, "Objectives", Inches(0.6), Inches(0.4), Inches(11), Inches(0.7),
             font=FONT_HEAD, size=32, bold=True, color=TEXT_PRIMARY)
    add_text(s, "Four concrete questions, one experimental program",
             Inches(0.6), Inches(1.05), Inches(11), Inches(0.4),
             font=FONT_BODY, size=13, italic=True, color=TEXT_SECONDARY)

    objectives = [
        ("01", "Compare codecs", "MP3 vs AAC vs Opus across 24–256 kbps", GOOGLE_BLUE_HEX),
        ("02", "Find transparency", "Bitrate at which listeners can no longer tell", GOOGLE_RED_HEX),
        ("03", "Build the player", "Web Audio API + HRTF, deployed and live", GOOGLE_YELLOW_HEX),
        ("04", "Stress the network", "Loss / jitter → spatial breakdown point", GOOGLE_GREEN_HEX),
    ]
    for i, (num, head, body, color) in enumerate(objectives):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.2)
        y = Inches(1.7 + row * 2.4)
        neumorphic_card(s, x, y, Inches(6.0), Inches(2.2))
        # number badge
        n = slide_num_circle(s, x + Inches(0.35), y + Inches(0.35), Inches(0.9), color, num)
        add_text(s, head, x + Inches(1.5), y + Inches(0.45), Inches(4.5), Inches(0.55),
                 font=FONT_HEAD, size=20, bold=True, color=TEXT_PRIMARY)
        add_text(s, body, x + Inches(1.5), Inches(y.inches + 1.0), Inches(4.5), Inches(1.0),
                 font=FONT_BODY, size=14, color=TEXT_SECONDARY)
    add_speaker_notes(s,
        "We organised the project around four concrete objectives. "
        "First, compare MP3, AAC and Opus across a wide bitrate ladder. "
        "Second, find each codec's transparency point - the bitrate at which a listener can no "
        "longer tell the compressed signal from the original. "
        "Third, build a working spatial audio player on the open Web Audio API and validate the "
        "binaural cues it produces. "
        "Fourth, stress the network - packet loss and jitter - to find when spatial immersion breaks. "
        "[25 seconds]"
    )
    add_footer(s, 3, TOTAL)

    # ===== SLIDE 4 — System & Methodology =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s)
    add_text(s, "System & Methodology", Inches(0.6), Inches(0.4), Inches(11), Inches(0.7),
             font=FONT_HEAD, size=32, bold=True, color=TEXT_PRIMARY)
    add_text(s, "End-to-end: encode → analyse → spatialise → stress",
             Inches(0.6), Inches(1.05), Inches(11), Inches(0.4),
             font=FONT_BODY, size=13, italic=True, color=TEXT_SECONDARY)

    # Pipeline steps (4 horizontal cards)
    steps = [
        ("FFmpeg", "libfdk-aac · libmp3lame · libopus", GOOGLE_BLUE_HEX),
        ("Metrics", "SNR + PESQ-derived ODG", GOOGLE_RED_HEX),
        ("Spatial player", "Web Audio API + HRTF", GOOGLE_YELLOW_HEX),
        ("Stress sim", "Packet loss + jitter sliders", GOOGLE_GREEN_HEX),
    ]
    for i, (h, body, color) in enumerate(steps):
        x = Inches(0.6 + i * 3.12)
        neumorphic_card(s, x, Inches(1.8), Inches(2.92), Inches(2.5))
        add_circle(s, x + Inches(0.25), Inches(2.05), Inches(0.5), color)
        add_text(s, h, x + Inches(0.25), Inches(2.7), Inches(2.6), Inches(0.45),
                 font=FONT_HEAD, size=16, bold=True, color=TEXT_PRIMARY)
        add_text(s, body, x + Inches(0.25), Inches(3.15), Inches(2.6), Inches(1.0),
                 font=FONT_BODY, size=12, color=TEXT_SECONDARY)

    # ABX protocol callout
    neumorphic_card(s, Inches(0.6), Inches(4.6), Inches(6.0), Inches(2.1),
                    fill_hex="EFF4FB", accent_hex=GOOGLE_BLUE_HEX)
    add_text(s, "ABX listening test", Inches(0.95), Inches(4.75), Inches(5.5), Inches(0.5),
             font=FONT_HEAD, size=18, bold=True, color=GOOGLE_BLUE)
    add_text(s,
             "Forced choice: A vs B vs X.\n"
             "3 codecs × 5 bitrates × 5 reps = 75 trials.\n"
             "Web interface; CSV log per participant.",
             Inches(0.95), Inches(5.25), Inches(5.5), Inches(1.4),
             font=FONT_BODY, size=13, color=TEXT_SECONDARY)

    neumorphic_card(s, Inches(7.0), Inches(4.6), Inches(5.7), Inches(2.1),
                    fill_hex="F4F0EF", accent_hex=GOOGLE_RED_HEX)
    add_text(s, "Stress design", Inches(7.35), Inches(4.75), Inches(5.5), Inches(0.5),
             font=FONT_HEAD, size=18, bold=True, color=GOOGLE_RED)
    add_text(s,
             "Loss: 0–20 % · Jitter: 0–200 ms.\n"
             "Mono / Stereo / HRTF compared.\n"
             "Metric = Immersion Degradation Index.",
             Inches(7.35), Inches(5.25), Inches(5.5), Inches(1.4),
             font=FONT_BODY, size=13, color=TEXT_SECONDARY)
    add_speaker_notes(s,
        "Our pipeline is straightforward. "
        "FFmpeg with libfdk-aac, libmp3lame and libopus produces the encoded fixtures. "
        "We compute SNR and a PESQ-derived ODG for objective quality. "
        "We render binaural audio in the browser through the Web Audio API's HRTF panner. "
        "And we simulate packet loss and jitter directly in the audio graph. "
        "Subjective evaluation is a forced-choice ABX protocol with 75 trials per participant. "
        "Stress testing compares mono, stereo and binaural HRTF under matched bitrates. "
        "[40 seconds]"
    )
    add_footer(s, 4, TOTAL)

    # ===== SLIDE 5 — Codec Comparison Results =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s)
    add_text(s, "Codec Comparison Results", Inches(0.6), Inches(0.4), Inches(11), Inches(0.7),
             font=FONT_HEAD, size=32, bold=True, color=TEXT_PRIMARY)
    add_text(s, "Rate–distortion: SNR misleads, ODG tells the truth",
             Inches(0.6), Inches(1.05), Inches(11), Inches(0.4),
             font=FONT_BODY, size=13, italic=True, color=TEXT_SECONDARY)

    # Left: figure card (placeholder/image)
    left_card = neumorphic_card(s, Inches(0.6), Inches(1.7), Inches(7.0), Inches(5.0))
    add_text(s,
             "[ Figure — Rate–Distortion (ODG) ]\n"
             "drop ./report/figures/rd_curve_odg.png",
             Inches(0.9), Inches(3.8), Inches(6.4), Inches(1.0),
             font=FONT_BODY, size=11, italic=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_text(s, "Figure 5 — ODG vs. bitrate (PESQ-derived)",
             Inches(0.9), Inches(6.05), Inches(6.4), Inches(0.4),
             font=FONT_BODY, size=11, italic=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # If we have the image, we'd add it here. Since we don't, the placeholder text is shown above.
    rd_path = FIG_DIR / "rd_curve_odg.png"
    if rd_path.exists():
        s.shapes.add_picture(str(rd_path), Inches(0.9), Inches(2.0), width=Inches(6.4))

    # Right: stats column
    stats = [
        ("48 kbps", "Opus transparency (pilot)", GOOGLE_GREEN_HEX),
        ("64 kbps", "MP3 / AAC transparency", GOOGLE_BLUE_HEX),
        ("≈ 20 ms", "Opus algorithmic delay", GOOGLE_YELLOW_HEX),
    ]
    for i, (big, label, color) in enumerate(stats):
        y = Inches(1.7 + i * 1.6)
        neumorphic_card(s, Inches(7.85), y, Inches(4.85), Inches(1.5))
        # color bar
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(7.85), y, Inches(0.18), Inches(1.5))
        bar.adjustments[0] = 0.5
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(color)
        _set_no_outline(bar)
        add_text(s, big, Inches(8.2), y + Inches(0.2), Inches(4.5), Inches(0.85),
                 font=FONT_HEAD, size=32, bold=True, color=TEXT_PRIMARY)
        add_text(s, label, Inches(8.2), y + Inches(0.95), Inches(4.5), Inches(0.45),
                 font=FONT_BODY, size=12, color=TEXT_SECONDARY)

    add_speaker_notes(s,
        "The rate-distortion curves split the ranking. "
        "By raw signal-to-noise ratio, Opus dominates. "
        "But SNR understates AAC because AAC deliberately puts its quantisation noise where "
        "perception masks it. "
        "When we move to a perceptually weighted ODG, the picture clarifies: Opus reaches "
        "transparency near 48 kbps in pilot listening; MP3 and AAC need closer to 64. "
        "Combined with Opus's roughly 20-millisecond algorithmic delay, this makes Opus the "
        "obvious choice for real-time spatial audio. [40 seconds]"
    )
    add_footer(s, 5, TOTAL)

    # ===== SLIDE 6 — ABX & Transparency =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s)
    add_text(s, "ABX Test & Transparency", Inches(0.6), Inches(0.4), Inches(11), Inches(0.7),
             font=FONT_HEAD, size=32, bold=True, color=TEXT_PRIMARY)
    add_text(s, "Pilot data — full cohort study in progress",
             Inches(0.6), Inches(1.05), Inches(11), Inches(0.4),
             font=FONT_BODY, size=13, italic=True, color=TEXT_SECONDARY)

    # Big result strip
    neumorphic_card(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.4),
                    fill_hex="EEF1F6")
    rows = [
        ("MP3", "≈ 64 kbps", "54 % (≈ chance)", "MOS 4.3", GOOGLE_BLUE_HEX),
        ("AAC", "≈ 64 kbps", "53 % (≈ chance)", "MOS 4.4", GOOGLE_RED_HEX),
        ("Opus", "≈ 48 kbps", "52 % (≈ chance)", "MOS 4.5", GOOGLE_GREEN_HEX),
    ]
    # Header row
    hdr_y = Inches(1.85)
    cols_x = [Inches(1.0), Inches(3.6), Inches(6.4), Inches(9.4)]
    headers = ["Codec", "Transparency bitrate", "ABX accuracy at threshold", "MOS at threshold"]
    for x, txt in zip(cols_x, headers):
        add_text(s, txt, x, hdr_y, Inches(3.0), Inches(0.4),
                 font=FONT_HEAD, size=12, bold=True, color=TEXT_DIM)
    # Data rows
    for i, (codec, br, acc, mos, color) in enumerate(rows):
        y = Inches(2.35 + i * 0.55)
        # Codec pill
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y - Inches(0.05),
                                   Inches(1.0), Inches(0.5))
        pill.adjustments[0] = 0.5
        pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor.from_string(color)
        _set_no_outline(pill)
        tf = pill.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = codec
        run.font.name = FONT_HEAD; run.font.size = Pt(13); run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for x, val in zip(cols_x[1:], (br, acc, mos)):
            add_text(s, val, x, y, Inches(3.0), Inches(0.5),
                     font=FONT_BODY, size=14, color=TEXT_PRIMARY,
                     anchor=MSO_ANCHOR.MIDDLE)

    # Bottom takeaway
    neumorphic_card(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.3),
                    fill_hex="EFF4FB", accent_hex=GOOGLE_BLUE_HEX)
    add_text(s, "Takeaway",
             Inches(1.0), Inches(4.55), Inches(11), Inches(0.5),
             font=FONT_HEAD, size=12, bold=True, color=GOOGLE_BLUE)
    add_text(s,
             "Opus reaches transparency roughly 1.3× faster than MP3/AAC.\n"
             "All three converge on near-chance discrimination at 64 kbps and above.",
             Inches(1.0), Inches(5.0), Inches(11.5), Inches(1.7),
             font=FONT_HEAD, size=18, color=TEXT_PRIMARY)

    add_speaker_notes(s,
        "The pilot ABX results align with the objective metrics. "
        "Opus reaches near-chance discrimination - our operational definition of transparency - at "
        "around 48 kilobits per second. MP3 and AAC need about 64. "
        "The MOS approximations cluster between 4.3 and 4.5 at threshold, all comfortably above the "
        "acceptable quality line of 3.5. "
        "These numbers are pilot-scale; the protocol, web interface and analysis notebook are all "
        "deployed - the full participant study is the next step. [35 seconds]"
    )
    add_footer(s, 6, TOTAL)

    # ===== SLIDE 7 — Network Stress =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s)
    add_text(s, "Network Stress: Spatial Breaks First",
             Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             font=FONT_HEAD, size=32, bold=True, color=TEXT_PRIMARY)
    add_text(s, "HRTF binaural streams are roughly 2× more fragile than stereo",
             Inches(0.6), Inches(1.05), Inches(11), Inches(0.4),
             font=FONT_BODY, size=13, italic=True, color=TEXT_SECONDARY)

    # Three big stat cards
    stats = [
        ("> 20 %", "Mono — packet loss tolerated", GOOGLE_BLUE_HEX),
        ("≈ 15 %", "Stereo — packet loss tolerated", GOOGLE_YELLOW_HEX),
        ("≈ 10 %", "Spatial (HRTF) — breakdown", GOOGLE_RED_HEX),
    ]
    for i, (big, label, color) in enumerate(stats):
        x = Inches(0.6 + i * 4.15)
        neumorphic_card(s, x, Inches(1.75), Inches(3.95), Inches(2.4))
        # Top color bar
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, Inches(1.75), Inches(3.95), Inches(0.18))
        bar.adjustments[0] = 0.5
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(color)
        _set_no_outline(bar)
        add_text(s, big, x + Inches(0.3), Inches(2.05), Inches(3.6), Inches(1.1),
                 font=FONT_HEAD, size=44, bold=True, color=TEXT_PRIMARY)
        add_text(s, label, x + Inches(0.3), Inches(3.25), Inches(3.6), Inches(0.7),
                 font=FONT_BODY, size=13, color=TEXT_SECONDARY)

    # Mechanism explainer
    neumorphic_card(s, Inches(0.6), Inches(4.45), Inches(12.1), Inches(2.25),
                    fill_hex="EFF4FB", accent_hex=GOOGLE_RED_HEX)
    add_text(s, "Why?",
             Inches(1.0), Inches(4.6), Inches(11), Inches(0.5),
             font=FONT_HEAD, size=12, bold=True, color=GOOGLE_RED)
    add_text(s,
             "HRTF rendering relies on phase-coherent two-channel delivery.\n"
             "A single late or dropped packet disrupts ITD/ILD and momentarily collapses the scene.",
             Inches(1.0), Inches(5.05), Inches(11.5), Inches(1.7),
             font=FONT_HEAD, size=17, color=TEXT_PRIMARY)
    add_speaker_notes(s,
        "Network impairments hit spatial audio harder than they hit mono or stereo. "
        "Mono tolerates 20% packet loss before quality collapses. "
        "Stereo holds out to about 15. "
        "Binaural HRTF breaks down by 10. "
        "The reason is mechanical: HRTF rendering depends on phase-coherent delivery of both "
        "channels. A single late or dropped packet disrupts the inter-aural cues and temporarily "
        "collapses the spatial scene. This is exactly the failure mode we need to engineer against. "
        "[30 seconds]"
    )
    add_footer(s, 7, TOTAL)

    # ===== SLIDE 8 — LIVE DEMO =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s, color_hex="DEE3EB")
    # Big neumorphic demo card
    neumorphic_card(s, Inches(0.7), Inches(0.7), Inches(11.9), Inches(6.1),
                    fill_hex="EEF1F6")
    # Big LIVE label
    add_pill(s, "LIVE DEMO", Inches(1.2), Inches(1.2), fill_hex=GOOGLE_RED_HEX, size=14)
    add_text(s, "Spatial Audio Lab",
             Inches(1.2), Inches(1.95), Inches(11), Inches(1.0),
             font=FONT_HEAD, size=44, bold=True, color=TEXT_PRIMARY)
    add_text(s,
             "vanowarna.com/spatial-audio-and-high-fidelity-streaming",
             Inches(1.2), Inches(2.85), Inches(11), Inches(0.5),
             font=FONT_BODY, size=14, italic=True, color=GOOGLE_BLUE)
    # Demo bullets as cards
    demo_steps = [
        ("Drop a sound source", "Drag the orb in 3D space — listener stays centred", GOOGLE_BLUE_HEX),
        ("Switch codecs live", "Original ↔ MP3 ↔ AAC ↔ Opus mid-playback", GOOGLE_GREEN_HEX),
        ("Inject network stress", "Loss & jitter sliders — feel the spatial collapse", GOOGLE_RED_HEX),
    ]
    for i, (h, body, color) in enumerate(demo_steps):
        y = Inches(3.6 + i * 1.0)
        neumorphic_card(s, Inches(1.2), y, Inches(11.0), Inches(0.85),
                        fill_hex="F2F5FA")
        add_circle(s, Inches(1.4), y + Inches(0.18), Inches(0.5), color)
        # number inside circle could be added but keep clean
        add_text(s, h, Inches(2.1), y + Inches(0.1), Inches(4.5), Inches(0.5),
                 font=FONT_HEAD, size=14, bold=True, color=TEXT_PRIMARY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, body, Inches(2.1), y + Inches(0.42), Inches(8.8), Inches(0.45),
                 font=FONT_BODY, size=11, color=TEXT_SECONDARY)
    add_speaker_notes(s,
        "Now the live demonstration. "
        "I'll open the Spatial Audio Lab in the browser. "
        "First I'll add a source and drag it through the listener's space - notice how the perceived "
        "position changes smoothly because the panner uses HRTF. "
        "Next I'll switch codec mid-playback - Original to MP3 to AAC to Opus - so you can hear "
        "the differences in real time without losing position. "
        "Finally I'll push the loss slider up - you'll hear the spatial scene start to collapse "
        "well before the audio itself sounds 'broken'. That's the breakdown point we measured. "
        "[75 seconds]"
    )
    add_footer(s, 8, TOTAL)

    # ===== SLIDE 9 — Conclusions & Future Work =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s)
    add_text(s, "Conclusions", Inches(0.6), Inches(0.4), Inches(11), Inches(0.7),
             font=FONT_HEAD, size=32, bold=True, color=TEXT_PRIMARY)
    add_text(s, "Pick Opus, ship spatial, design for fragility",
             Inches(0.6), Inches(1.05), Inches(11), Inches(0.4),
             font=FONT_BODY, size=13, italic=True, color=TEXT_SECONDARY)

    # Findings card
    neumorphic_card(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.0))
    add_text(s, "What we found", Inches(0.95), Inches(1.85), Inches(5.5), Inches(0.5),
             font=FONT_HEAD, size=18, bold=True, color=GOOGLE_BLUE)
    findings = [
        ("Opus wins", "Best quality / latency / bitrate balance"),
        ("Spatial fragility", "HRTF breaks at ~10 % loss; stereo holds longer"),
        ("Bandwidth myth", "5G has the bytes — jitter & loss are the real risk"),
    ]
    for i, (h, b) in enumerate(findings):
        y = Inches(2.45 + i * 1.4)
        add_circle(s, Inches(0.95), y + Inches(0.08), Inches(0.4),
                   [GOOGLE_BLUE_HEX, GOOGLE_RED_HEX, GOOGLE_GREEN_HEX][i])
        add_text(s, h, Inches(1.55), y, Inches(4.8), Inches(0.45),
                 font=FONT_HEAD, size=14, bold=True, color=TEXT_PRIMARY)
        add_text(s, b, Inches(1.55), y + Inches(0.5), Inches(4.8), Inches(0.7),
                 font=FONT_BODY, size=12, color=TEXT_SECONDARY)

    # Future work card
    neumorphic_card(s, Inches(6.85), Inches(1.7), Inches(5.85), Inches(5.0),
                    fill_hex="F4F0EF", accent_hex=GOOGLE_RED_HEX)
    add_text(s, "Where this goes next",
             Inches(7.2), Inches(1.85), Inches(5.5), Inches(0.5),
             font=FONT_HEAD, size=18, bold=True, color=GOOGLE_RED)
    futures = [
        "MPEG-H 3D Object-Based Audio",
        "Personalised HRTFs (CNN-based)",
        "Opus multistream for ambisonics",
        "Neural codecs (EnCodec / SoundStream)",
        "WebRTC pipeline on real 4G/5G traces",
    ]
    for i, f in enumerate(futures):
        y = Inches(2.45 + i * 0.75)
        add_circle(s, Inches(7.2), y + Inches(0.05), Inches(0.25), GOOGLE_RED_HEX)
        add_text(s, f, Inches(7.6), y, Inches(5), Inches(0.5),
                 font=FONT_BODY, size=14, color=TEXT_PRIMARY,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_speaker_notes(s,
        "To wrap up. "
        "Opus is the codec to deploy for real-time spatial audio - it gives the best ratio of "
        "quality to bits and the lowest delay. "
        "Spatial streams are inherently more fragile than stereo, so production systems must engineer "
        "against jitter and burst loss. "
        "On 5G, bandwidth is no longer the bottleneck; coherence is. "
        "Future work: object-based audio with MPEG-H, personalised HRTFs, Opus multistream for "
        "ambisonics, neural codecs, and a real network evaluation through WebRTC. [25 seconds]"
    )
    add_footer(s, 9, TOTAL)

    # ===== SLIDE 10 — Thank You =====
    s = prs.slides.add_slide(blank)
    style_slide_background(s, color_hex="DEE3EB")
    neumorphic_card(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.6),
                    fill_hex="EEF1F6")
    add_text(s, "Thank you", Inches(1.8), Inches(2.1), Inches(10), Inches(1.2),
             font=FONT_HEAD, size=64, bold=True, color=TEXT_PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, "Questions?", Inches(1.8), Inches(3.4), Inches(10), Inches(0.6),
             font=FONT_BODY, size=20, italic=True, color=TEXT_SECONDARY,
             align=PP_ALIGN.CENTER)
    # Google color row
    cy = Inches(4.5)
    add_circle(s, Inches(5.7), cy, Inches(0.4), GOOGLE_BLUE_HEX)
    add_circle(s, Inches(6.2), cy, Inches(0.4), GOOGLE_RED_HEX)
    add_circle(s, Inches(6.7), cy, Inches(0.4), GOOGLE_YELLOW_HEX)
    add_circle(s, Inches(7.2), cy, Inches(0.4), GOOGLE_GREEN_HEX)
    add_text(s, "vanowarna.com/spatial-audio-and-high-fidelity-streaming",
             Inches(1.8), Inches(5.3), Inches(10), Inches(0.5),
             font=FONT_BODY, size=14, italic=True, color=GOOGLE_BLUE,
             align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Thank you. We have time for questions. The full report and live player are linked above. "
        "[5 seconds]"
    )
    add_footer(s, 10, TOTAL)

    prs.save(str(OUT_PPTX))
    print(f"Saved: {OUT_PPTX}")


def slide_num_circle(slide, x, y, d, color_hex, num):
    """Circle with a number inside."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    _set_no_outline(c)
    _add_outer_shadow(c, dist=8000, blur=30000, alpha=30000, color=color_hex)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = num
    run.font.name = FONT_HEAD; run.font.size = Pt(16); run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return c


if __name__ == "__main__":
    build()
